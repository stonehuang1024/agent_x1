"""
PowerPoint Tools Module - Presentation creation and manipulation.

Provides tools for:
- Creating PowerPoint presentations from structured content
- Reading existing presentations
- Adding slides, text, images
- Exporting to PDF via LibreOffice
"""

import logging
from pathlib import Path
from typing import Dict, Any, List, Optional

from ..core.tool import Tool

logger = logging.getLogger(__name__)


def create_presentation(
    slides: List[Dict[str, Any]],
    output_path: str,
    title: str = "Presentation",
    theme: str = "default"
) -> Dict[str, Any]:
    """
    Create a PowerPoint presentation from structured slide data.

    Each slide dict supports:
      - title (str): Slide title
      - content (str | list): Body text or bullet list
      - layout (str): 'title', 'content', 'two_column', 'blank' (default: 'content')
      - image_path (str): Optional image path to embed
      - notes (str): Speaker notes

    Args:
        slides: List of slide dicts
        output_path: Output .pptx file path
        title: Presentation title (for metadata)
        theme: Color theme name (default, dark, blue)

    Returns:
        Dictionary with result
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    THEMES = {
        "default": {"bg": RGBColor(0xFF, 0xFF, 0xFF), "title_fg": RGBColor(0x1F, 0x49, 0x7D), "body_fg": RGBColor(0x00, 0x00, 0x00)},
        "dark": {"bg": RGBColor(0x1E, 0x1E, 0x2E), "title_fg": RGBColor(0xCB, 0xD2, 0xEA), "body_fg": RGBColor(0xB0, 0xB8, 0xC8)},
        "blue": {"bg": RGBColor(0x00, 0x33, 0x66), "title_fg": RGBColor(0xFF, 0xFF, 0xFF), "body_fg": RGBColor(0xCC, 0xE5, 0xFF)},
    }
    colors = THEMES.get(theme, THEMES["default"])

    try:
        prs = Presentation()
        prs.core_properties.title = title

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        LAYOUT_BLANK = 6

        created_slides = []

        for i, slide_data in enumerate(slides):
            slide_layout = prs.slide_layouts[LAYOUT_BLANK]
            slide = prs.slides.add_slide(slide_layout)

            # Background color
            from pptx.oxml.ns import qn
            from lxml import etree
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = colors["bg"]

            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            notes_text = slide_data.get("notes", "")
            image_path = slide_data.get("image_path")

            # Title text box
            if slide_title:
                txBox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(1.2)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = slide_title
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = colors["title_fg"]

            # Content text box
            body_top = Inches(1.7)
            body_height = slide_height - Inches(2.2)
            if image_path:
                body_height = slide_height - Inches(2.2) - Inches(2.5)

            if content:
                txBox2 = slide.shapes.add_textbox(
                    Inches(0.5), body_top,
                    slide_width - Inches(1), body_height
                )
                tf2 = txBox2.text_frame
                tf2.word_wrap = True

                if isinstance(content, list):
                    for j, bullet in enumerate(content):
                        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                        p.level = 0
                        run = p.add_run()
                        run.text = f"• {bullet}"
                        run.font.size = Pt(18)
                        run.font.color.rgb = colors["body_fg"]
                else:
                    p = tf2.paragraphs[0]
                    run = p.add_run()
                    run.text = str(content)
                    run.font.size = Pt(18)
                    run.font.color.rgb = colors["body_fg"]

            # Image
            if image_path:
                img_resolved = Path(image_path).expanduser().resolve()
                if img_resolved.exists():
                    img_top = slide_height - Inches(2.8)
                    slide.shapes.add_picture(
                        str(img_resolved),
                        Inches(0.5), img_top,
                        slide_width - Inches(1), Inches(2.5)
                    )

            # Speaker notes
            if notes_text:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = notes_text

            created_slides.append({
                "index": i + 1,
                "title": slide_title,
                "has_image": bool(image_path),
                "has_notes": bool(notes_text)
            })

        out = Path(output_path).expanduser().resolve()
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))

        return {
            "output_path": str(out),
            "slide_count": len(slides),
            "slides": created_slides,
            "file_size_bytes": out.stat().st_size,
            "success": True
        }
    except Exception as e:
        logger.exception("[CreatePresentation] Failed")
        return {"error": str(e)}


def read_presentation(path: str) -> Dict[str, Any]:
    """
    Read content from an existing PowerPoint file.

    Args:
        path: Path to .pptx file

    Returns:
        Dictionary with slide content and metadata
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    try:
        resolved = Path(path).expanduser().resolve()
        if not resolved.exists():
            return {"error": f"File not found: {path}"}

        prs = Presentation(str(resolved))
        slides_data = []

        for i, slide in enumerate(prs.slides):
            texts = []
            images = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
                if shape.shape_type == 13:
                    images.append(shape.name)

            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides_data.append({
                "index": i + 1,
                "texts": texts,
                "image_count": len(images),
                "notes": notes
            })

        return {
            "path": str(resolved),
            "slide_count": len(prs.slides),
            "title": prs.core_properties.title,
            "author": prs.core_properties.author,
            "slides": slides_data
        }
    except Exception as e:
        logger.exception("[ReadPresentation] Failed")
        return {"error": str(e), "path": path}


def add_slide(
    path: str,
    title: str,
    content: str,
    position: Optional[int] = None
) -> Dict[str, Any]:
    """
    Add a new slide to an existing presentation.

    Args:
        path: Path to existing .pptx file
        title: Slide title
        content: Slide body text
        position: Insert position (1-indexed). Default: append at end.

    Returns:
        Dictionary with result
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    try:
        resolved = Path(path).expanduser().resolve()
        if not resolved.exists():
            return {"error": f"File not found: {path}"}

        prs = Presentation(str(resolved))
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        if title:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
            tf = txBox.text_frame
            tf.text = title
            tf.paragraphs[0].runs[0].font.size = Pt(28)
            tf.paragraphs[0].runs[0].font.bold = True

        if content:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(4.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            tf2.text = content
            tf2.paragraphs[0].runs[0].font.size = Pt(18)

        if position is not None:
            from pptx.oxml.ns import qn
            slides_element = prs.slides._sldIdLst
            moved = slides_element[-1]
            slides_element.remove(moved)
            insert_idx = max(0, min(position - 1, len(slides_element)))
            slides_element.insert(insert_idx, moved)

        prs.save(str(resolved))
        return {
            "path": str(resolved),
            "total_slides": len(prs.slides),
            "added_title": title,
            "success": True
        }
    except Exception as e:
        logger.exception("[AddSlide] Failed")
        return {"error": str(e)}


def export_presentation_to_pdf(pptx_path: str, output_path: Optional[str] = None) -> Dict[str, Any]:
    """
    Export a PowerPoint file to PDF using LibreOffice.

    Requires LibreOffice to be installed (libreoffice command).

    Args:
        pptx_path: Path to .pptx file
        output_path: Output PDF path. Default: same dir as pptx with .pdf extension.

    Returns:
        Dictionary with result
    """
    import subprocess

    try:
        resolved = Path(pptx_path).expanduser().resolve()
        if not resolved.exists():
            return {"error": f"File not found: {pptx_path}"}

        out = Path(output_path).expanduser().resolve() if output_path else resolved.with_suffix(".pdf")
        out.parent.mkdir(parents=True, exist_ok=True)

        result = subprocess.run(
            [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(out.parent), str(resolved)
            ],
            capture_output=True, text=True, timeout=120
        )

        if result.returncode != 0:
            return {
                "error": f"LibreOffice conversion failed: {result.stderr}",
                "pptx_path": str(resolved)
            }

        generated = resolved.with_suffix(".pdf")
        if generated.exists() and str(generated) != str(out):
            generated.rename(out)

        return {
            "pptx_path": str(resolved),
            "pdf_path": str(out),
            "file_size_bytes": out.stat().st_size if out.exists() else None,
            "success": True
        }
    except FileNotFoundError:
        return {"error": "LibreOffice not found. Install libreoffice to use this feature."}
    except Exception as e:
        logger.exception("[ExportPPTtoPDF] Failed")
        return {"error": str(e)}


# Tool Definitions
CREATE_PRESENTATION_TOOL = Tool(
    name="create_presentation",
    description=(
        "Create a PowerPoint (.pptx) presentation from a list of slides. "
        "Each slide has: title (str), content (str or list of bullets), "
        "layout, image_path (optional), notes (optional). "
        "Themes: default, dark, blue."
    ),
    parameters={
        "type": "object",
        "properties": {
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"description": "Body text string or list of bullet strings"},
                        "image_path": {"type": "string"},
                        "notes": {"type": "string"}
                    }
                },
                "description": "List of slide dicts"
            },
            "output_path": {"type": "string", "description": "Output .pptx file path"},
            "title": {"type": "string", "description": "Presentation title (metadata)"},
            "theme": {"type": "string", "enum": ["default", "dark", "blue"], "description": "Color theme"}
        },
        "required": ["slides", "output_path"]
    },
    func=create_presentation,
    timeout_seconds=60,
    max_output_chars=5000,
)

READ_PRESENTATION_TOOL = Tool(
    name="read_presentation",
    description="Read text content and metadata from an existing .pptx PowerPoint file.",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to .pptx file"}
        },
        "required": ["path"]
    },
    func=read_presentation,
    timeout_seconds=30,
    max_output_chars=30000,
    is_readonly=True,
)

ADD_SLIDE_TOOL = Tool(
    name="add_slide",
    description="Add a new slide to an existing .pptx presentation. Optionally specify insert position.",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to existing .pptx file"},
            "title": {"type": "string", "description": "Slide title"},
            "content": {"type": "string", "description": "Slide body text"},
            "position": {"type": "integer", "description": "Insert position (1-indexed). Default: append."}
        },
        "required": ["path", "title", "content"]
    },
    func=add_slide,
    timeout_seconds=30,
    max_output_chars=5000,
)

EXPORT_PPT_TO_PDF_TOOL = Tool(
    name="export_presentation_to_pdf",
    description=(
        "Convert a .pptx presentation to PDF using LibreOffice. "
        "Requires LibreOffice to be installed on the system."
    ),
    parameters={
        "type": "object",
        "properties": {
            "pptx_path": {"type": "string", "description": "Path to .pptx file"},
            "output_path": {"type": "string", "description": "Output PDF path (default: same dir as pptx)"}
        },
        "required": ["pptx_path"]
    },
    func=export_presentation_to_pdf,
    timeout_seconds=120,
    max_output_chars=5000,
)

PPT_TOOLS = [
    CREATE_PRESENTATION_TOOL,
    READ_PRESENTATION_TOOL,
    ADD_SLIDE_TOOL,
    EXPORT_PPT_TO_PDF_TOOL,
]
